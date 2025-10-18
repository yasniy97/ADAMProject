import streamlit as st
import os
import json
import hashlib
import requests
from datetime import datetime
from typing import List, Dict
import PyPDF2
import openpyxl
from pptx import Presentation
import io
from bs4 import BeautifulSoup

# Page configuration
st.set_page_config(
    page_title="A.D.A.M - Agile Digital Assistance for Managers",
    page_icon="🤖",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Custom CSS for ChatGPT-like styling
st.markdown("""
<style>
    .main {
        background-color: #343541;
    }
    .stTextInput > div > div > input {
        background-color: #40414f;
        color: #ececf1;
        border: 1px solid #565869;
        border-radius: 5px;
    }
    .stButton > button {
        background-color: #10a37f;
        color: white;
        border-radius: 5px;
        border: none;
        padding: 0.5rem 1rem;
    }
    .stButton > button:hover {
        background-color: #0d8c6d;
    }
    .chat-message {
        padding: 1.5rem;
        border-radius: 0.5rem;
        margin-bottom: 1rem;
        display: flex;
        flex-direction: column;
    }
    .user-message {
        background-color: #343541;
        color: #ececf1;
    }
    .assistant-message {
        background-color: #444654;
        color: #ececf1;
    }
    .sidebar .sidebar-content {
        background-color: #202123;
    }
    .follow-up-questions {
        background-color: #2d2e3a;
        padding: 1rem;
        border-radius: 0.5rem;
        margin-top: 1rem;
    }
    .follow-up-btn {
        background-color: #40414f;
        color: #ececf1;
        padding: 0.5rem 1rem;
        margin: 0.3rem;
        border-radius: 5px;
        border: 1px solid #565869;
        cursor: pointer;
    }
    .mode-button-active {
        background-color: #10a37f !important;
        color: white !important;
        border: 2px solid #0d8c6d !important;
    }
    .mode-button-inactive {
        background-color: #40414f !important;
        color: #8e8ea0 !important;
        border: 1px solid #565869 !important;
    }
    .search-result-card {
        background-color: #40414f;
        padding: 1.5rem;
        border-radius: 0.5rem;
        margin-bottom: 1rem;
        border-left: 4px solid #10a37f;
    }
    .search-result-title {
        color: #10a37f;
        font-size: 1.2rem;
        font-weight: bold;
        margin-bottom: 0.5rem;
    }
    .search-result-url {
        color: #8e8ea0;
        font-size: 0.9rem;
        margin-bottom: 0.5rem;
    }
    .search-result-content {
        color: #ececf1;
        line-height: 1.6;
    }
</style>
""", unsafe_allow_html=True)

# Initialize session state
def init_session_state():
    if 'authenticated' not in st.session_state:
        st.session_state.authenticated = False
    if 'username' not in st.session_state:
        st.session_state.username = None
    if 'chat_history' not in st.session_state:
        st.session_state.chat_history = []
    if 'conversations' not in st.session_state:
        st.session_state.conversations = []
    if 'current_conversation_id' not in st.session_state:
        st.session_state.current_conversation_id = None
    if 'llm_type' not in st.session_state:
        st.session_state.llm_type = 'ollama_cloud'
    if 'llm_provider' not in st.session_state:
        st.session_state.llm_provider = 'ollama_cloud'
    if 'api_key' not in st.session_state:
        st.session_state.api_key = ''
    if 'search_mode' not in st.session_state:
        st.session_state.search_mode = False
    if 'search_results' not in st.session_state:
        st.session_state.search_results = []

# User authentication functions
def hash_password(password: str) -> str:
    return hashlib.sha256(password.encode()).hexdigest()

def load_users():
    if os.path.exists('users.json'):
        with open('users.json', 'r') as f:
            return json.load(f)
    return {}

def save_users(users):
    with open('users.json', 'w') as f:
        json.dump(users, f)

def create_user(username: str, password: str):
    users = load_users()
    if username in users:
        return False
    users[username] = hash_password(password)
    save_users(users)
    return True

def verify_user(username: str, password: str) -> bool:
    users = load_users()
    if username not in users:
        return False
    return users[username] == hash_password(password)

def reset_password(username: str, new_password: str) -> bool:
    users = load_users()
    if username not in users:
        return False
    users[username] = hash_password(new_password)
    save_users(users)
    return True

# Document processing functions
def extract_text_from_pdf(file) -> str:
    pdf_reader = PyPDF2.PdfReader(file)
    text = ""
    for page in pdf_reader.pages:
        text += page.extract_text()
    return text

def extract_text_from_pptx(file) -> str:
    prs = Presentation(file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_text_from_excel(file) -> str:
    wb = openpyxl.load_workbook(file)
    text = ""
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        text += f"\n--- Sheet: {sheet} ---\n"
        for row in ws.iter_rows(values_only=True):
            text += " | ".join([str(cell) if cell is not None else "" for cell in row]) + "\n"
    return text

def process_uploaded_files(files) -> str:
    combined_text = ""
    for file in files:
        file_extension = file.name.split('.')[-1].lower()
        try:
            if file_extension == 'pdf':
                combined_text += f"\n--- {file.name} ---\n" + extract_text_from_pdf(file)
            elif file_extension == 'pptx' or file_extension == 'ppt':
                combined_text += f"\n--- {file.name} ---\n" + extract_text_from_pptx(file)
            elif file_extension == 'xlsx' or file_extension == 'xls':
                combined_text += f"\n--- {file.name} ---\n" + extract_text_from_excel(file)
        except Exception as e:
            st.error(f"Error processing {file.name}: {str(e)}")
    return combined_text

# Check if query is PM-related
def is_pm_related(query: str) -> bool:
    pm_keywords = [
        'project management', 'program management', 'portfolio management',
        'agile', 'scrum', 'sprint', 'kanban', 'sdlc', 'software engineering',
        'pmbok', 'pmi', 'prince2', 'waterfall', 'stakeholder', 'risk management',
        'project manager', 'backlog', 'standup', 'retrospective', 'velocity',
        'burndown', 'epic', 'user story', 'product owner', 'scrum master'
    ]
    query_lower = query.lower()
    return any(keyword in query_lower for keyword in pm_keywords)

# Web scraping for PM resources
def search_pm_resources(query: str) -> str:
    """Search PM-specific resources"""
    resources = []
    pm_sites = {
        'pmi.org': 'PMI - Project Management Institute',
        'projectmanagement.com': 'ProjectManagement.com',
        'prince2.com': 'PRINCE2'
    }
    
    # Note: In production, implement actual web scraping or API calls
    # This is a placeholder for the structure
    resources.append(f"Searching PM resources for: {query}")
    resources.append("Note: Integrate actual web scraping or API calls to pmi.org, projectmanagement.com, and prince2.com")
    
    return "\n".join(resources)

# Web search function
def ollama_web_search(query: str, api_key: str, max_results: int = 5) -> Dict:
    """
    Search the web using Ollama's web search API
    """
    try:
        url = "https://ollama.com/api/web_search"
        headers = {
            'Authorization': f'Bearer {api_key}',
            'Content-Type': 'application/json'
        }
        payload = {
            'query': query,
            'max_results': min(max_results, 10)  # Ensure max is 10
        }
        
        response = requests.post(url, headers=headers, json=payload)
        response.raise_for_status()
        
        return response.json()
    except Exception as e:
        return {
            'error': f"Error performing web search: {str(e)}",
            'results': []
        }

def format_search_results(search_data: Dict) -> str:
    """Format search results for display using the correct API response format"""
    if 'error' in search_data:
        return f"⚠️ {search_data['error']}"
    
    results = search_data.get('results', [])
    if not results:
        return "No results found."
    
    formatted = ""
    for idx, result in enumerate(results, 1):
        title = result.get('title', 'No title')
        url = result.get('url', '#')
        content = result.get('content', 'No content available')
        
        formatted += f"""
<div class='search-result-card'>
    <div class='search-result-title'>{idx}. {title}</div>
    <div class='search-result-url'>🔗 {url}</div>
    <div class='search-result-content'>{content}</div>
</div>
"""
    
    return formatted

# LLM interaction functions
def call_openai_stream(prompt: str, api_key: str, model: str = "gpt-4"):
    """Call OpenAI API with streaming"""
    try:
        import openai
        openai.api_key = api_key
        
        stream = openai.ChatCompletion.create(
            model=model,
            messages=[{"role": "user", "content": prompt}],
            stream=True
        )
        
        for chunk in stream:
            if chunk.choices[0].delta.get("content"):
                yield chunk.choices[0].delta.content
    except Exception as e:
        yield f"Error calling OpenAI: {str(e)}"

def call_anthropic_stream(prompt: str, api_key: str, model: str = "claude-sonnet-4-5-20250929"):
    """Call Anthropic API with streaming"""
    try:
        import anthropic
        
        client = anthropic.Anthropic(api_key=api_key)
        
        with client.messages.stream(
            model=model,
            max_tokens=4096,
            messages=[{"role": "user", "content": prompt}]
        ) as stream:
            for text in stream.text_stream:
                yield text
    except Exception as e:
        yield f"Error calling Anthropic: {str(e)}"

def call_ollama_local_stream(prompt: str, model: str = "deepseek-r1"):
    """Call local Ollama instance with streaming"""
    try:
        import ollama
        
        # Add no_think parameter for deepseek-r1
        options = {}
        if model == "deepseek-r1":
            options['no_think'] = True
        
        stream = ollama.chat(
            model=model,
            messages=[{'role': 'user', 'content': prompt}],
            options=options,
            stream=True
        )
        
        for chunk in stream:
            yield chunk['message']['content']
    except Exception as e:
        yield f"Error calling local LLM: {str(e)}"

def call_ollama_cloud_stream(prompt: str, api_key: str, model: str = "deepseek-v3.1:671b-cloud"):
    """Call cloud Ollama instance with streaming"""
    try:
        from ollama import Client
        
        client = Client(
            host="https://ollama.com",
            headers={'Authorization': f'Bearer {api_key}'}
        )
        
        messages = [{'role': 'user', 'content': prompt}]
        
        for part in client.chat(model, messages=messages, stream=True):
            yield part['message']['content']
    except Exception as e:
        yield f"Error calling cloud LLM: {str(e)}"

def generate_response_stream(query: str, context: str = "", llm_provider: str = "ollama_cloud", api_key: str = "", model: str = ""):
    """Generate streaming response using appropriate LLM"""
    
    # Build the prompt
    prompt = f"You are A.D.A.M (Agile Digital Assistance for Managers), a specialized assistant for project management topics.\n\n"
    
    if context:
        prompt += f"Context from uploaded documents:\n{context}\n\n"
    
    # Check if PM-related and add PM resource context
    if is_pm_related(query):
        pm_context = search_pm_resources(query)
        prompt += f"PM Resource Context:\n{pm_context}\n\n"
    
    prompt += f"User Query: {query}\n\nProvide a comprehensive, well-formatted response."
    
    # Call appropriate LLM with streaming
    if not api_key:
        yield "Error: API key required. Please enter your API key in the sidebar."
        return
    
    if llm_provider == "ollama_cloud":
        for chunk in call_ollama_cloud_stream(prompt, api_key, model or "deepseek-v3.1:671b-cloud"):
            yield chunk
    elif llm_provider == "openai":
        for chunk in call_openai_stream(prompt, api_key, model or "gpt-4"):
            yield chunk
    elif llm_provider == "anthropic":
        for chunk in call_anthropic_stream(prompt, api_key, model or "claude-sonnet-4-5-20250929"):
            yield chunk
    else:
        for chunk in call_ollama_local_stream(prompt, model or "deepseek-r1"):
            yield chunk

def generate_follow_up_questions(query: str, response: str) -> List[str]:
    """Generate follow-up questions based on the conversation"""
    # This could be enhanced with LLM generation
    # For now, returning template questions
    if is_pm_related(query):
        return [
            "Can you provide more details about best practices?",
            "What are common challenges in implementing this?",
            "How does this relate to other PM methodologies?"
        ]
    else:
        return [
            "Can you elaborate on that?",
            "What are the key considerations?",
            "Are there any examples or case studies?"
        ]

# Login/Registration page
def login_page():
    st.markdown("""
    <div style='text-align: center; padding: 2rem;'>
        <h1 style='color: #10a37f;'>🤖 A.D.A.M</h1>
        <h3 style='color: #ececf1;'>Agile Digital Assistance for Managers</h3>
    </div>
    """, unsafe_allow_html=True)
    
    tab1, tab2, tab3 = st.tabs(["Login", "Register", "Reset Password"])
    
    with tab1:
        st.subheader("Login")
        username = st.text_input("Username", key="login_username")
        password = st.text_input("Password", type="password", key="login_password")
        
        if st.button("Login", key="login_btn"):
            if verify_user(username, password):
                st.session_state.authenticated = True
                st.session_state.username = username
                st.rerun()
            else:
                st.error("Invalid username or password")
    
    with tab2:
        st.subheader("Create New Account")
        new_username = st.text_input("Username", key="reg_username")
        new_password = st.text_input("Password", type="password", key="reg_password")
        confirm_password = st.text_input("Confirm Password", type="password", key="reg_confirm")
        
        if st.button("Register", key="reg_btn"):
            if new_password != confirm_password:
                st.error("Passwords do not match")
            elif len(new_password) < 6:
                st.error("Password must be at least 6 characters")
            elif create_user(new_username, new_password):
                st.success("Account created successfully! Please login.")
            else:
                st.error("Username already exists")
    
    with tab3:
        st.subheader("Reset Password")
        reset_username = st.text_input("Username", key="reset_username")
        new_pass = st.text_input("New Password", type="password", key="reset_new_pass")
        confirm_new_pass = st.text_input("Confirm New Password", type="password", key="reset_confirm")
        
        if st.button("Reset Password", key="reset_btn"):
            if new_pass != confirm_new_pass:
                st.error("Passwords do not match")
            elif len(new_pass) < 6:
                st.error("Password must be at least 6 characters")
            elif reset_password(reset_username, new_pass):
                st.success("Password reset successfully!")
            else:
                st.error("Username not found")

# Main chat interface
def main_page():
    # Sidebar
    with st.sidebar:
        st.markdown("""
        <div style='text-align: center; padding: 1rem;'>
            <h2 style='color: #10a37f;'>🤖 A.D.A.M</h2>
            <p style='color: #ececf1; font-size: 0.8rem;'>Agile Digital Assistance for Managers</p>
        </div>
        """, unsafe_allow_html=True)
        
        st.markdown(f"**User:** {st.session_state.username}")
        
        # LLM Selection
        st.markdown("---")
        st.subheader("LLM Settings")
        
        st.session_state.llm_provider = st.selectbox(
            "Select LLM Provider",
            ["ollama_cloud", "openai", "anthropic"],
            format_func=lambda x: {
                "ollama_cloud": "Ollama Cloud (DeepSeek)",
                "openai": "OpenAI (GPT-4)",
                "anthropic": "Anthropic (Claude)"
            }[x]
        )
        
        st.session_state.api_key = st.text_input(
            "API Key",
            type="password",
            value=st.session_state.api_key,
            help=f"Required API key for {st.session_state.llm_provider}"
        )
        
        # Show model selection based on provider
        if st.session_state.llm_provider == "openai":
            st.session_state.selected_model = st.selectbox(
                "Model",
                ["gpt-4", "gpt-4-turbo", "gpt-3.5-turbo"]
            )
        elif st.session_state.llm_provider == "anthropic":
            st.session_state.selected_model = st.selectbox(
                "Model",
                ["claude-sonnet-4-5-20250929", "claude-opus-4-1-20250514", "claude-sonnet-4-20250514"]
            )
        else:
            st.session_state.selected_model = "deepseek-v3.1:671b-cloud"
        
        # Chat History
        st.markdown("---")
        st.subheader("Mode Selection")
        
        # Search Mode Toggle with custom styling
        col1, col2 = st.columns(2)
        with col1:
            chat_class = "mode-button-active" if not st.session_state.search_mode else "mode-button-inactive"
            st.markdown(f"""
                <style>
                    div[data-testid="column"]:nth-of-type(1) button {{
                        {f'background-color: #10a37f !important; border: 2px solid #0d8c6d !important;' if not st.session_state.search_mode else 'background-color: #40414f !important; border: 1px solid #565869 !important;'}
                    }}
                </style>
            """, unsafe_allow_html=True)
            if st.button("💬 Chat", use_container_width=True, key="chat_mode_btn"):
                st.session_state.search_mode = False
                st.rerun()
        
        with col2:
            search_class = "mode-button-active" if st.session_state.search_mode else "mode-button-inactive"
            st.markdown(f"""
                <style>
                    div[data-testid="column"]:nth-of-type(2) button {{
                        {f'background-color: #10a37f !important; border: 2px solid #0d8c6d !important;' if st.session_state.search_mode else 'background-color: #40414f !important; border: 1px solid #565869 !important;'}
                    }}
                </style>
            """, unsafe_allow_html=True)
            if st.button("🔍 Search", use_container_width=True, key="search_mode_btn"):
                st.session_state.search_mode = True
                st.rerun()
        
        # Display current mode
        mode_text = "🔍 Search Mode Active" if st.session_state.search_mode else "💬 Chat Mode Active"
        mode_color = "#10a37f" if st.session_state.search_mode else "#10a37f"
        st.markdown(f"<p style='text-align: center; color: {mode_color}; font-weight: bold;'>{mode_text}</p>", unsafe_allow_html=True)
        
        st.markdown("---")
        st.subheader("Chat History")
        
        if st.button("➕ New Chat", use_container_width=True):
            st.session_state.chat_history = []
            st.session_state.current_conversation_id = None
            st.session_state.search_results = []
            st.rerun()
        
        if st.button("🗑️ Clear History", use_container_width=True):
            st.session_state.chat_history = []
            st.session_state.conversations = []
            st.session_state.current_conversation_id = None
            st.session_state.search_results = []
            st.rerun()
        
        # Display saved conversations
        for idx, conv in enumerate(st.session_state.conversations):
            if st.button(f"💬 {conv['title'][:30]}...", key=f"conv_{idx}", use_container_width=True):
                st.session_state.chat_history = conv['messages']
                st.session_state.current_conversation_id = conv['id']
                st.rerun()
        
        st.markdown("---")
        if st.button("🚪 Logout", use_container_width=True):
            st.session_state.authenticated = False
            st.session_state.username = None
            st.rerun()
    
    # Main chat area
    if st.session_state.search_mode:
        st.markdown("<h1 style='color: #ececf1;'>🔍 A.D.A.M Search</h1>", unsafe_allow_html=True)
        st.info("Search Mode: Find information on any topic using web search")
    else:
        st.markdown("<h1 style='color: #ececf1;'>💬 A.D.A.M Chat</h1>", unsafe_allow_html=True)
    
    # File upload
    uploaded_files = st.file_uploader(
        "Upload documents (PDF, PPT, Excel, MPP)",
        type=['pdf', 'pptx', 'ppt', 'xlsx', 'xls', 'mpp'],
        accept_multiple_files=True
    )
    
    document_context = ""
    if uploaded_files:
        with st.spinner("Processing documents..."):
            document_context = process_uploaded_files(uploaded_files)
        st.success(f"Processed {len(uploaded_files)} document(s)")
    
    # Display chat history
    chat_container = st.container()
    with chat_container:
        for msg_idx, message in enumerate(st.session_state.chat_history):
            if message['role'] == 'user':
                # Check if it's a search query
                if message.get('search_query', False):
                    st.markdown(f"""
                    <div class='chat-message user-message'>
                        <strong>🔍 You searched for:</strong><br>{message['content'].replace('🔍 Search: ', '')}
                    </div>
                    """, unsafe_allow_html=True)
                else:
                    st.markdown(f"""
                    <div class='chat-message user-message'>
                        <strong>You:</strong><br>{message['content']}
                    </div>
                    """, unsafe_allow_html=True)
            else:
                # Check if it's search results
                if message.get('search_results', False):
                    st.markdown("""
                    <div class='chat-message assistant-message'>
                        <strong>A.D.A.M Search Results:</strong><br>
                    """, unsafe_allow_html=True)
                    st.markdown(message['content'], unsafe_allow_html=True)
                    st.markdown("</div>", unsafe_allow_html=True)
                else:
                    st.markdown(f"""
                    <div class='chat-message assistant-message'>
                        <strong>A.D.A.M:</strong><br>{message['content']}
                    </div>
                    """, unsafe_allow_html=True)
                
                # Display follow-up questions only for the last assistant message
                if 'follow_ups' in message and msg_idx == len(st.session_state.chat_history) - 1:
                    st.markdown("<div class='follow-up-questions'><strong>Follow-up questions:</strong></div>", unsafe_allow_html=True)
                    cols = st.columns(len(message['follow_ups']))
                    for idx, follow_up in enumerate(message['follow_ups']):
                        with cols[idx]:
                            if st.button(follow_up, key=f"followup_{message['timestamp']}_{idx}"):
                                st.session_state.pending_query = follow_up
                                st.rerun()
    
    # Chat input at the bottom
    st.markdown("<br>", unsafe_allow_html=True)
    
    # Handle pending query from follow-up
    if 'pending_query' in st.session_state:
        query = st.session_state.pending_query
        del st.session_state.pending_query
    else:
        query = st.chat_input("ADAM> Type your message here...")
    
    if query:
        # Add user message
        st.session_state.chat_history.append({
            'role': 'user',
            'content': query,
            'timestamp': datetime.now().isoformat()
        })
        
        # Create placeholder for streaming response
        response_placeholder = st.empty()
        full_response = ""
        
        # Stream the response
        with response_placeholder.container():
            st.markdown("""
            <div class='chat-message user-message'>
                <strong>You:</strong><br>""" + query + """
            </div>
            """, unsafe_allow_html=True)
            
            # Create assistant message container
            assistant_message = st.markdown("""
            <div class='chat-message assistant-message'>
                <strong>A.D.A.M:</strong><br>
            </div>
            """, unsafe_allow_html=True)
            
            # Stream the response
            for chunk in generate_response_stream(
                query,
                document_context,
                st.session_state.llm_provider,
                st.session_state.api_key,
                st.session_state.get('selected_model', '')
            ):
                full_response += chunk
                assistant_message.markdown(f"""
                <div class='chat-message assistant-message'>
                    <strong>A.D.A.M:</strong><br>{full_response}
                </div>
                """, unsafe_allow_html=True)
        
        # Generate follow-up questions after streaming completes
        follow_ups = generate_follow_up_questions(query, full_response)
        
        # Add assistant message to history
        st.session_state.chat_history.append({
            'role': 'assistant',
            'content': full_response,
            'follow_ups': follow_ups,
            'timestamp': datetime.now().isoformat()
        })
        
        # Save conversation
        if not st.session_state.current_conversation_id:
            conv_id = f"conv_{datetime.now().timestamp()}"
            st.session_state.conversations.append({
                'id': conv_id,
                'title': query[:50],
                'messages': st.session_state.chat_history.copy()
            })
            st.session_state.current_conversation_id = conv_id
        else:
            # Update existing conversation
            for conv in st.session_state.conversations:
                if conv['id'] == st.session_state.current_conversation_id:
                    conv['messages'] = st.session_state.chat_history.copy()
                    break
        
        st.rerun()

# Main app
def main():
    init_session_state()
    
    if not st.session_state.authenticated:
        login_page()
    else:
        main_page()

if __name__ == "__main__":
    main()